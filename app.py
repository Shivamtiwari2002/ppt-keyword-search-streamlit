# app.py
import streamlit as st
import pandas as pd
import os
import re
import zipfile
import tempfile
from io import BytesIO
from pptx import Presentation
from rapidfuzz import fuzz
import html
from pathlib import Path
from openai import OpenAI

# ---------------- CONFIG ----------------
OPENAI_KEY = None
if "OPENAI_API_KEY" in st.secrets:
    OPENAI_KEY = st.secrets["OPENAI_API_KEY"]
else:
    OPENAI_KEY = os.getenv("OPENAI_API_KEY")

if not OPENAI_KEY:
    st.warning("OpenAI API key not found. AI features will show an error until you set OPENAI_API_KEY.")
else:
    client = OpenAI(api_key=OPENAI_KEY)

# ---------------- PAGE CONFIG ----------------
st.set_page_config(page_title="Impact Analysis Tool", layout="wide")

# ---------------- UI THEME WITH CSS ----------------
st.markdown("""
<style>
body {background: linear-gradient(140deg, rgba(247,249,255,0.85) 0%, rgba(237,242,255,0.85) 40%, rgba(255,255,255,0.85) 100%); font-family: 'Segoe UI', sans-serif;}
.header-box {background: linear-gradient(135deg, #0047FF, #3F8CFF); padding: 36px; border-radius: 20px; text-align: center; color: white; font-size: 36px; font-weight: 800; margin-bottom: 35px; box-shadow: 0px 8px 25px rgba(0,40,140,0.25);}
.section-card {background: white; padding: 26px; border-radius: 18px; border: 1px solid #DFE6FF; box-shadow: 0px 10px 22px rgba(0,60,160,0.08); margin-bottom: 28px;}
.stFileUploader>div>div {border: 2px dashed #0047FF !important; background: #EFF3FF !important; border-radius: 16px !important;}
input[type="text"] {border: 2px solid #0047FF !important; border-radius: 14px !important; padding: 14px !important; background: #EFF3FF !important; font-size: 16px !important; color: #0033CC !important; font-weight: 600 !important;}
.stButton>button {background: linear-gradient(135deg, #0047FF, #2F6BFF) !important; color: white !important; border-radius: 12px !important; padding: 12px 28px !important; font-size: 17px !important; border: none !important; font-weight: 600 !important; box-shadow: 0px 5px 15px rgba(0,0,0,0.17); transition: 0.2s ease-in-out;}
.stButton>button:hover {transform: translateY(-2px); box-shadow: 0px 7px 18px rgba(0,0,0,0.22);}
.table-container {width: 100%; border-collapse: collapse; margin-top: 22px;}
.table-container th {background: #0047FF; color: white; padding: 12px; text-align: left;}
.table-container tr:nth-child(even) {background: #F0F4FF;} .table-container tr:hover {background: #E5EBFF;} .table-container td {padding: 12px; border-bottom: 1px solid #D0D8FF; font-size: 15px;}
.footer {text-align: center; margin-top: 40px; padding: 14px; font-size: 14px; font-weight: 600; color: #0047FF;}
.mark {background: #FFF176;}
</style>
""", unsafe_allow_html=True)

st.markdown("<div class='header-box'>Impact Analysis Tool</div>", unsafe_allow_html=True)

# ---------------- Helper: PPT → HTML / text extraction ----------------
def ppt_to_html_slides(file_path):
    prs = Presentation(file_path)
    slides_out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""
        raw_text_pieces = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                raw_text_pieces.append(shape.text)
                txt = html.escape(shape.text).replace("\n", "<br>")
                parts.append(f"<p>{txt}</p>")
        html_content = "<div><h3>Slide %d</h3><h4>%s</h4>%s</div>" % (i, html.escape(title), "".join(parts))
        slides_out.append({
            "slide_no": i,
            "title": title,
            "html": html_content,
            "raw_text": " ".join(raw_text_pieces)
        })
    return slides_out

# ---------------- Search helpers ----------------
def highlight_terms(html_text, keyword):
    pattern = re.compile(re.escape(keyword), re.IGNORECASE)
    return pattern.sub(lambda m: f"<mark class='mark'>{m.group(0)}</mark>", html_text)

def search_slides(slides, keyword, mode="exact_phrase", threshold=80):
    results = []
    for s in slides:
        text_for_search = s["raw_text"] or ""
        lowered = text_for_search.lower()
        if mode == "exact_phrase" and keyword.lower() in lowered:
            results.append({**s, "score": 100})
        elif mode == "exact" and keyword.lower() in lowered:
            results.append({**s, "score": 100})
        elif mode == "fuzzy":
            score = fuzz.partial_ratio(keyword.lower(), lowered)
            if score >= threshold:
                results.append({**s, "score": score})
    return results

def extract_zip_pptx(zip_file_path):
    temp_dir = tempfile.mkdtemp()
    with zipfile.ZipFile(zip_file_path, "r") as z:
        z.extractall(temp_dir)
    pptx_files = []
    for root, _, files in os.walk(temp_dir):
        for f in files:
            if f.lower().endswith(".pptx"):
                pptx_files.append(os.path.join(root, f))
    return pptx_files

# ---------------- SIDEBAR ----------------
with st.sidebar:
    st.markdown("### Search Mode")
    search_mode = st.radio("", ["exact_phrase (recommended)", "exact", "fuzzy"], index=0)
    threshold = st.slider("Fuzzy threshold", 60, 100, 85)
    st.markdown("---")
    st.markdown("Made by SKT")

# ---------------- UPLOAD ----------------
st.markdown("<div class='section-card'>", unsafe_allow_html=True)
st.markdown("### Upload PPTX / ZIP Files")
uploaded_files = st.file_uploader("", type=["pptx", "zip"], accept_multiple_files=True)
st.markdown("</div>", unsafe_allow_html=True)

# ---------------- PROCESS UPLOADED FILES ----------------
pptx_paths = []
all_slides_text = []

if uploaded_files:
    upload_dir = Path(tempfile.gettempdir()) / "ppt_uploads"
    upload_dir.mkdir(parents=True, exist_ok=True)
    for uf in uploaded_files:
        tmp_path = upload_dir / uf.name
        with open(tmp_path, "wb") as f:
            f.write(uf.getbuffer())
        if uf.name.lower().endswith(".pptx"):
            pptx_paths.append(str(tmp_path))
        else:
            pptx_paths.extend(extract_zip_pptx(str(tmp_path)))

    for p in pptx_paths:
        try:
            slides = ppt_to_html_slides(p)
            for s in slides:
                all_slides_text.append({
                    "file": os.path.basename(p),
                    "slide_no": s["slide_no"],
                    "title": s["title"],
                    "text": s["raw_text"]
                })
        except Exception as e:
            st.error(f"Error processing {p}: {e}")

# ---------------- AI Retrieval & Chat ----------------
def retrieve_relevant_slides(question, slides, top_n=5, min_score=20):
    ranked = []
    q = question.lower()
    for s in slides:
        txt = (s.get("text") or "").lower()
        score = fuzz.partial_ratio(q, txt)
        ranked.append((score, s))
    ranked.sort(key=lambda x: x[0], reverse=True)
    results = [s for score, s in ranked[:top_n] if score >= min_score]
    return results

def ask_ai_question(user_question, context_text):
    if OPENAI_KEY is None:
        return "OpenAI API key not configured. Set OPENAI_API_KEY in st.secrets or environment."

    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are an assistant that answers questions using ONLY the provided PPT slide content. "
                        "If the answer is not present in the context, say you couldn't find it. "
                        "Always cite slide numbers when possible in square brackets like [Slide 3]."
                    )
                },
                {
                    "role": "user",
                    "content": f"Context:\n{context_text}\n\nQuestion:\n{user_question}\n\nAnswer concisely and cite slide numbers used."
                }
            ],
            temperature=0.2,
            max_tokens=700
        )
        return response.choices[0].message.content

    except Exception as e:
        error_str = str(e)
        if "insufficient_quota" in error_str or "429" in error_str:
            return (
                "⚠️ AI request failed: You have exceeded your OpenAI quota. "
                "Please check your plan, billing, or reduce the number of slides / tokens."
            )
        return f"AI Error: {error_str}"

# ---------------- UI: AI Chat ----------------
st.markdown("<div class='section-card'>", unsafe_allow_html=True)
st.markdown("### 🤖 AI Chat — Two Modes (Search-based & Full Chatbot)")

col1, col2 = st.columns([2, 1])
with col1:
    chat_scope = st.radio("Chat scope", ["Search-based (use keyword search results)", "Full Chatbot (use entire PPT)"], index=1)
    user_question = st.text_input("Ask anything about the PPT:", "", placeholder="e.g., In which visualization is Job XYZ used?")
with col2:
    top_k = st.number_input("Top-K slides to use for context", value=5, min_value=1, max_value=10)
    min_score = st.slider("Min relevance score for retrieval", 0, 100, 20)
    st.markdown("Use **Search-based** when you want the AI to rely on your keyword search results. Use **Full Chatbot** to search entire PPT.")
ask_ai_btn = st.button("✨ Ask AI")
st.markdown("</div>", unsafe_allow_html=True)

if ask_ai_btn:
    if not uploaded_files or not all_slides_text:
        st.error("Please upload PPTX / ZIP files first.")
    elif not user_question.strip():
        st.error("Please enter a question.")
    else:
        with st.spinner("Retrieving relevant slides and asking AI..."):
            if chat_scope.startswith("Search-based"):
                last_search_results = st.session_state.get("last_search_matches", [])
                if not last_search_results:
                    relevant = retrieve_relevant_slides(user_question, all_slides_text, top_n=top_k, min_score=min_score)
                else:
                    mapped = []
                    for m in last_search_results:
                        for s in all_slides_text:
                            if s["file"] == m["File"] and s["slide_no"] == int(m["Slide"]):
                                mapped.append(s)
                    relevant = mapped[:top_k] if mapped else retrieve_relevant_slides(user_question, all_slides_text, top_n=top_k, min_score=min_score)
            else:
                relevant = retrieve_relevant_slides(user_question, all_slides_text, top_n=top_k, min_score=min_score)

            if not relevant:
                st.warning("No sufficiently relevant slides found.")
            else:
                context_text = ""
                for s in relevant:
                    context_text += f"[Slide {s['slide_no']} — {s['file']}]\nTitle: {s['title']}\n{s['text']}\n\n"
                ai_answer = ask_ai_question(user_question, context_text)

                st.markdown("<div class='section-card'>", unsafe_allow_html=True)
                st.markdown("### 🤖 AI Answer")
                st.write(ai_answer)
                st.markdown("</div>", unsafe_allow_html=True)

                st.markdown("### 📌 Source Slides Used")
                for s in relevant:
                    st.markdown(f"**Slide {s['slide_no']} — {s['title']}**")
                    st.write(s["text"][:600] + ("..." if len(s["text"]) > 600 else ""))

# ---------------- KEYWORD SEARCH ----------------
st.markdown("<div class='section-card'>", unsafe_allow_html=True)
keyword = st.text_input("Enter Keyword", "", placeholder="e.g. PSD Manager")
st.markdown("</div>", unsafe_allow_html=True)

search_btn = st.button("🔍 Search")
results_all = []

if search_btn:
    if not pptx_paths:
        st.error("Please upload PPTX / ZIP files.")
    elif not keyword.strip():
        st.error("Please enter a keyword.")
    else:
        mode_clean = search_mode.split()[0]
        with st.spinner("Searching slides…"):
            for p in pptx_paths:
                try:
                    slides = ppt_to_html_slides(p)
                    matches = search_slides(slides, keyword, mode_clean, threshold)
                    for m in matches:
                        highlighted = highlight_terms(m["html"], keyword)
                        results_all.append({
                            "File": os.path.basename(p),
                            "Slide": m["slide_no"],
                            "Title": m["title"],
                            "Score": m["score"],
                            "HTML": highlighted,
                            "raw_text": m["raw_text"]
                        })
                except Exception as e:
                    st.error(f"Error reading {p}: {e}")
        st.session_state["last_search_matches"] = results_all.copy()
        st.success(f"{len(results_all)} matches found.")

# ---------------- SHOW KEYWORD SEARCH RESULTS (NO PREVIEW) ----------------
if results_all:
    df = pd.DataFrame(results_all).drop(columns=["HTML", "raw_text"])
    st.markdown("<div class='section-card'>", unsafe_allow_html=True)
    st.markdown("### Search Results")

    def render_table(df):
        table = "<table class='table-container'>"
        table += "<tr>" + "".join(f"<th>{c}</th>" for c in df.columns) + "</tr>"
        for _, row in df.iterrows():
            table += "<tr>" + "".join(f"<td>{html.escape(str(x))}</td>" for x in row) + "</tr>"
        table += "</table>"
        return table

    st.markdown(render_table(df), unsafe_allow_html=True)

    excel_buffer = BytesIO()
    with pd.ExcelWriter(excel_buffer, engine="openpyxl") as writer:
        df.to_excel(writer, index=False)
    excel_buffer.seek(0)
    st.download_button("⬇ Download Results (Excel)", excel_buffer.getvalue(), "ppt_search_results.xlsx")

    st.markdown("</div>", unsafe_allow_html=True)

# ---------------- FOOTER ----------------
st.markdown("<div class='footer'>Made by SKT</div>", unsafe_allow_html=True)
